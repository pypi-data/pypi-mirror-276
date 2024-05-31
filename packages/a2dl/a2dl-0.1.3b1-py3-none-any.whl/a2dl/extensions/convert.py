import os.path
import time
from a2dl.core.constants import logger
from a2dl.core.library import Diaglibrary


def printlibrary(library):
    print('LIBRARY: ', library.name, library.icons)
    for icon in library.icons:
        print('_ICON: ', icon.name, icon.image)
        for variable in icon.variables:
            print('__TEXT: ', variable['name'], variable['content'])


def to_pptx(library: Diaglibrary, target: str, whitelist: list[str] = None, tpl: str = None):
    """
    Converts a Library into a Power Point Presentation with each Icon on a Slide
    :param library: a a2dl library object
    :param target: full filepath for the pptx to be created or overwritten
    :param whitelist: the title strings from asciidoc titles or names from library - icons attrributes which should get implemented
    :param tpl: a full path to a power point template
    :return:
    """
    import pptx.util

    # create new empty Presentation
    if not tpl:
        prs = pptx.Presentation()
    else:
        prs = pptx.Presentation(os.path.abspath(tpl))

    SLD_LAYOUT_TITLE = 0
    SLD_LAYOUT_SLIDE = 6  # 6 = blank

    # title slide
    title_slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = library.name
    subtitle.text = f'{str(time.gmtime().tm_year)}.{str(time.gmtime().tm_mon)}.{str(time.gmtime().tm_mday)}'

    # content slides
    slide_layout = prs.slide_layouts[SLD_LAYOUT_SLIDE]
    for icon in library.icons:
        slide = prs.slides.add_slide(slide_layout)

        # title
        h_width = prs.slide_width
        h_height = prs.slide_height * 0.1
        h_left = 0
        h_top = 0

        h_txBox = slide.shapes.add_textbox(h_left, h_top, h_width, h_height)  # Adding Shape object (Text Box)
        h_tf = h_txBox.text_frame
        h_tf.auto_size = 2  # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        p = h_tf.add_paragraph()
        p.text = icon.name
        p.font.bold = True

        # details
        t_width = prs.slide_width * 0.5
        t_height = prs.slide_height
        t_left = prs.slide_width * 0.5
        t_top = prs.slide_height * 0.1

        txBox = slide.shapes.add_textbox(t_left, t_top, t_width, t_height)  # Adding Shape object (Text Box)
        tf = txBox.text_frame
        tf.auto_size = 2  # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = True
        for variable in icon.variables:
            if whitelist:
                if variable['name'] in whitelist:
                    p = tf.add_paragraph()
                    p.text = variable['name'] + ':'
                    p.font.bold = True
                    p = tf.add_paragraph()
                    # todo: tigabeatz
                    try:
                        p.text = variable['content'][0]
                    except IndexError:
                        pass
            else:
                p = tf.add_paragraph()
                p.text = variable['name'] + ':'
                p.font.bold = True
                p = tf.add_paragraph()
                # todo: tigabeatz
                try:
                    p.text = variable['content'][0]
                except IndexError:
                    pass
        # image
        try:
            p_width = prs.slide_width * 0.5
            p_left = 0
            p_top = prs.slide_height * 0.1
            # todo: tigabeatz
            try:
                pic = slide.shapes.add_picture(icon.image, p_left, p_top, p_width)
            except AttributeError:

                pass

        except FileNotFoundError:
            pass
    try:
        prs.save(os.path.abspath(target))
    except FileNotFoundError as err:
        logger.error(f'Can not write file. Does the Directory exist? {err}')


def to_docx(library: Diaglibrary, target: str, whitelist: list[str] = None, tpl: str = None):
    """
    Converts a Library into a Word Document with each Icon on a Page
    :param library: a a2dl library object
    :param target: full filepath for the docx to be created or overwritten
    :param whitelist: the title strings from asciidoc titles or names from library - icons attrributes which should get implemented
    :param tpl: a full path to a  word template
    :return:
    """
    # printlibrary(library)
    from docx import Document
    from docx.shared import Inches

    if not tpl:
        document = Document()
    else:
        document = Document(os.path.abspath(tpl))

    document.add_heading(library.name, 0)

    document.add_paragraph(f'{str(time.gmtime().tm_year)}.{str(time.gmtime().tm_mon)}.{str(time.gmtime().tm_mday)}')

    for icon in library.icons:
        document.add_page_break()
        document.add_heading(icon.name, level=1)
        try:
            try:
                document.add_picture(icon.image, width=Inches(1.25))
            except AttributeError:
                pass
        except FileNotFoundError:
            pass
        for variable in icon.variables:
            if whitelist:
                if variable['name'] in whitelist:
                    document.add_heading(variable['name'], level=2)
                    try:
                        document.add_paragraph(variable['content'][0], style='Intense Quote')
                    except IndexError:
                        pass
            else:
                document.add_heading(variable['name'], level=2)
                try:
                    document.add_paragraph(variable['content'][0], style='Intense Quote')
                except IndexError:
                    pass
    try:
        document.save(os.path.abspath(target))
    except FileNotFoundError as err:
        logger.error(f'Can not write file. Does the Directory exist? {err}')


def to_ea(library: Diaglibrary, target: str, whitelist: list[str] = None, tpl: str = None):
    printlibrary(library)
    raise NotImplementedError
